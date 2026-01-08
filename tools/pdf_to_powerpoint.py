import fitz
import os
from pptx import Presentation
from pptx.util import Inches, Pt

def pdf_to_powerpoint(input_path, output_path):
    try:
        doc = fitz.open(input_path)
        prs = Presentation()
        
        # Define a blank slide layout (usually index 6 in default template)
        blank_slide_layout = prs.slide_layouts[6]
        
        for page_num, page in enumerate(doc):
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Extract text blocks
            blocks = page.get_text("blocks")
            for block in blocks:
                # fitz block: (x0, y0, x1, y1, "text", block_no, block_type)
                x0, y0, x1, y1, text, _, _ = block
                width = x1 - x0
                height = y1 - y0
                
                # Create text box roughly in position (converting points to Inches)
                # 1 Inch = 72 Points
                left = Inches(x0 / 72)
                top = Inches(y0 / 72)
                w = Inches(width / 72)
                h = Inches(height / 72)
                
                txBox = slide.shapes.add_textbox(left, top, w, h)
                tf = txBox.text_frame
                tf.text = text.strip()
            
            # Extract images
            image_list = page.get_images()
            for img_index, img in enumerate(image_list):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    ext = base_image["ext"]
                    
                    # Get image position on page
                    # This finds usages of the image on the current page
                    rects = page.get_image_rects(xref)
                    for rect in rects:
                        left = Inches(rect.x0 / 72)
                        top = Inches(rect.y0 / 72)
                        width = Inches(rect.width / 72)
                        height = Inches(rect.height / 72)
                        
                        # Save temp image to add to slide
                        import io
                        image_stream = io.BytesIO(image_bytes)
                        slide.shapes.add_picture(image_stream, left, top, width, height)
                except Exception:
                    continue

        prs.save(output_path)
        doc.close()
        
        return {'success': True, 'output_path': output_path, 'filename': os.path.basename(output_path)}
    except Exception as e:
        return {'success': False, 'error': str(e)}
