from pptx import Presentation
from reportlab.lib.pagesizes import landscape, letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch

def pptx_to_pdf(input_path, output_path):
    try:
        prs = Presentation(input_path)
        
        pdf = SimpleDocTemplate(output_path, pagesize=landscape(letter),
                               rightMargin=72, leftMargin=72,
                               topMargin=72, bottomMargin=72)
        
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'SlideTitle',
            parent=styles['Heading1'],
            fontSize=24,
            leading=28,
            spaceAfter=24,
            alignment=1
        )
        
        content_style = ParagraphStyle(
            'SlideContent',
            parent=styles['Normal'],
            fontSize=14,
            leading=18,
            spaceAfter=12,
            bulletIndent=20,
            leftIndent=20
        )
        
        story = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            story.append(Paragraph(f"Slide {slide_num}", styles['Heading3']))
            story.append(Spacer(1, 12))
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                            
                            if paragraph.level == 0:
                                story.append(Paragraph(text, title_style))
                            else:
                                bullet = '• ' if paragraph.level > 0 else ''
                                story.append(Paragraph(f"{bullet}{text}", content_style))
            
            if slide_num < len(prs.slides):
                story.append(PageBreak())
        
        if story:
            pdf.build(story)
        else:
            pdf.build([Paragraph("Empty presentation", styles['Normal'])])
        
        return {'success': True, 'output_path': output_path, 'filename': output_path.split('/')[-1]}
    except Exception as e:
        return {'success': False, 'error': str(e)}
